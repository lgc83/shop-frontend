from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def compact(text: str) -> str:
    return " ".join((text or "").split())


@dataclass(frozen=True)
class TextStyle:
    name: str | None
    size_pt: float | None
    bold: bool | None
    color: RGBColor | None


def first_run_style(shape) -> TextStyle:
    if not getattr(shape, "has_text_frame", False):
        return TextStyle(None, None, None, None)
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            f = r.font
            size_pt = float(f.size.pt) if f.size is not None else None
            color = None
            try:
                if f.color is not None and f.color.rgb is not None:
                    color = f.color.rgb
            except Exception:
                color = None
            return TextStyle(
                name=f.name,
                size_pt=size_pt,
                bold=f.bold,
                color=color,
            )
    return TextStyle(None, None, None, None)


def apply_style(run, style: TextStyle) -> None:
    f = run.font
    if style.name:
        f.name = style.name
    if style.size_pt:
        f.size = Pt(style.size_pt)
    if style.bold is not None:
        f.bold = style.bold
    if style.color is not None:
        f.color.rgb = style.color


def add_textbox(slide, left, top, width, height, text: str, style: TextStyle, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    apply_style(r, style)
    return box


def main() -> None:
    src = Path("이기창 포트폴리오(작업용).pptx")
    out = Path("이기창 포트폴리오_07쇼핑몰추가_수정본.pptx")

    prs = Presentation(str(src))

    # --- 1) TOC slide (index 3 => 0-based 2) add "07"
    toc = prs.slides[2]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Try to infer styles/positions from existing number/label shapes
    num_style = TextStyle(None, 36, True, RGBColor(255, 255, 255))
    label_style = TextStyle(None, 16, True, RGBColor(255, 255, 255))

    num_shapes = {}
    label_shapes = {}
    for sh in toc.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = compact(sh.text)
        if t in {"01", "02", "03", "04", "05", "06"}:
            num_shapes[t] = sh
        if t in {"프로젝트 개요", "업무흐름(ERP·MES)", "시스템 설계 개요", "시스템 아키텍처", "기술 스택 및 구현", "핵심 기능"}:
            label_shapes[t] = sh

    # Copy number style from "06" (best match to bottom row)
    if "06" in num_shapes:
        s = first_run_style(num_shapes["06"])
        num_style = TextStyle(s.name, s.size_pt or 36, True if s.bold is None else s.bold, s.color or RGBColor(255, 255, 255))

    # Copy label style from "핵심 기능" if present
    if "핵심 기능" in label_shapes:
        s = first_run_style(label_shapes["핵심 기능"])
        label_style = TextStyle(s.name, s.size_pt or 16, True if s.bold is None else s.bold, s.color or RGBColor(255, 255, 255))

    # Position: create a 3rd row centered below existing row 2
    bottom_row = [num_shapes.get("04"), num_shapes.get("05"), num_shapes.get("06")]
    bottoms = [sh.top + sh.height for sh in bottom_row if sh is not None]
    y_base = max(bottoms) if bottoms else int(slide_h * 0.68)

    num_w = Inches(1.2)
    num_h = Inches(0.55)
    label_w = Inches(3.0)
    label_h = Inches(0.4)

    y_num = y_base + Inches(0.25)
    y_label = y_num + Inches(0.55)

    # Clamp within slide
    max_y = slide_h - Inches(0.35)
    if y_label + label_h > max_y:
        y_num = slide_h - Inches(1.45)
        y_label = y_num + Inches(0.55)

    x_num = int((slide_w - num_w) / 2)
    x_label = int((slide_w - label_w) / 2)

    add_textbox(toc, x_num, y_num, num_w, num_h, "07", num_style, PP_ALIGN.CENTER)
    add_textbox(toc, x_label, y_label, label_w, label_h, "쇼핑몰(SHOP)", label_style, PP_ALIGN.CENTER)

    # --- 2) Add section cover slide for 07 based on slide 4's look (index 4 => 0-based 3)
    cover_ref = prs.slides[3]
    cover_layout = cover_ref.slide_layout
    cover = prs.slides.add_slide(cover_layout)

    # If layout doesn't include the actual typography, mimic using copied positions from ref slide
    # Find two key text boxes on the reference slide.
    ref_texts = {}
    for sh in cover_ref.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = compact(sh.text)
        if t in {"01", "프로젝트 개요"}:
            ref_texts[t] = sh

    # Fallback positions
    title_left = Inches(1.2)
    title_top = Inches(2.6)
    title_w = Inches(7.6)
    title_h = Inches(0.8)
    num_left = Inches(1.2)
    num_top = Inches(1.8)
    num_w2 = Inches(1.4)
    num_h2 = Inches(0.8)

    if "프로젝트 개요" in ref_texts:
        sh = ref_texts["프로젝트 개요"]
        title_left, title_top, title_w, title_h = sh.left, sh.top, sh.width, sh.height
        s = first_run_style(sh)
        label_style = TextStyle(s.name, s.size_pt or 28, True if s.bold is None else s.bold, s.color or RGBColor(255, 255, 255))
    else:
        label_style = TextStyle(label_style.name, 28, True, RGBColor(255, 255, 255))

    if "01" in ref_texts:
        sh = ref_texts["01"]
        num_left, num_top, num_w2, num_h2 = sh.left, sh.top, sh.width, sh.height
        s = first_run_style(sh)
        num_style = TextStyle(s.name, s.size_pt or 40, True if s.bold is None else s.bold, s.color or RGBColor(255, 255, 255))
    else:
        num_style = TextStyle(num_style.name, 40, True, RGBColor(255, 255, 255))

    add_textbox(cover, num_left, num_top, num_w2, num_h2, "07", num_style, PP_ALIGN.LEFT)
    add_textbox(cover, title_left, title_top, title_w, title_h, "쇼핑몰(SHOP)", label_style, PP_ALIGN.LEFT)

    # --- 3) Content slide for 07 (dark background, images, bullets)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    s07 = prs.slides.add_slide(blank_layout)

    # Dark background rectangle (full bleed)
    bg = s07.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(11, 18, 32)
    bg.line.color.rgb = RGBColor(11, 18, 32)

    # Title
    title_style = TextStyle(label_style.name, 30, True, RGBColor(255, 255, 255))
    add_textbox(
        s07,
        Inches(0.7),
        Inches(0.45),
        Inches(8.8),
        Inches(0.6),
        "07 쇼핑몰(SHOP) 기능/화면 개요",
        title_style,
        PP_ALIGN.LEFT,
    )

    # Bullets
    bullet_box = s07.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(4.5), Inches(3.9))
    tf = bullet_box.text_frame
    tf.clear()
    bullets = [
        "사용자: 상품 조회 · 주문 진입",
        "관리자: 상품 마스터 등록/수정/삭제",
        "판매상태(판매중/품절/비활성) 관리",
        "ERP·MES 연계 확장 고려(주문 → 재고/출고)",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)
        p.font.bold = False
        p.font.color.rgb = RGBColor(229, 231, 235)

    # Images (use existing project screenshots)
    img_main = Path("public/img/main.png")
    img_admin = Path("public/img/admin.png")

    img_w = Inches(4.2)
    img_h = Inches(2.35)
    x_img = Inches(5.2)
    y1 = Inches(1.25)
    y2 = Inches(3.85)

    if img_main.exists():
        s07.shapes.add_picture(str(img_main), x_img, y1, width=img_w, height=img_h)
        add_textbox(s07, x_img, y1 + img_h + Inches(0.08), img_w, Inches(0.25), "사용자 메인", TextStyle(None, 12, True, RGBColor(203, 213, 225)))
    if img_admin.exists():
        s07.shapes.add_picture(str(img_admin), x_img, y2, width=img_w, height=img_h)
        add_textbox(s07, x_img, y2 + img_h + Inches(0.08), img_w, Inches(0.25), "관리자 상품 마스터", TextStyle(None, 12, True, RGBColor(203, 213, 225)))

    prs.save(str(out))
    print("OK:", out.name)


if __name__ == "__main__":
    main()

