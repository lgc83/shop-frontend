from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def compact(text: str) -> str:
    return " ".join((text or "").split())


def main() -> None:
    pptx_path = Path("이기창 포트폴리오(작업용).pptx")
    prs = Presentation(str(pptx_path))
    print("file:", pptx_path)
    print("slides:", len(prs.slides))
    print("size:", int(prs.slide_width), int(prs.slide_height))

    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            t = compact(getattr(shape, "text", ""))
            if t:
                texts.append(t)
        preview = " | ".join(texts)
        if len(preview) > 220:
            preview = preview[:220] + "…"
        print(f"{idx:02d}:", preview)


if __name__ == "__main__":
    main()

