from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def compact(text: str) -> str:
    return " ".join((text or "").split())


def main() -> None:
    candidates = sorted(Path(".").glob("이기창 포트폴리오_*수정본.pptx"))
    if not candidates:
        raise SystemExit("no output pptx found")

    f = candidates[-1]
    prs = Presentation(str(f))

    toc = prs.slides[2]
    texts = []
    for sh in toc.shapes:
        if getattr(sh, "has_text_frame", False) and (sh.text or "").strip():
            texts.append(compact(sh.text))

    toc_has_07 = "07" in texts and any("쇼핑몰" in t for t in texts)

    last_texts = []
    for sh in prs.slides[-1].shapes:
        if getattr(sh, "has_text_frame", False) and (sh.text or "").strip():
            last_texts.append(compact(sh.text))

    print("file_ok")
    print("file:", f.name)
    print("slides:", len(prs.slides))
    print("toc_has_07:", toc_has_07)
    print("last_slide_text_preview:", " | ".join(last_texts)[:140])


if __name__ == "__main__":
    main()

